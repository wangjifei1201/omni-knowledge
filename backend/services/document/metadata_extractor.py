"""
Document Metadata Extractor Service
Uses LLM to extract document metadata from file content
"""

import io
import json
import re
from typing import Optional

from loguru import logger
from services.llm.llm import llm_service


class MetadataExtractor:
    """
    Extract document metadata using LLM.
    Supports batch extraction for multiple documents.
    """

    def __init__(self):
        self.max_preview_length = 2000  # Max chars to send to LLM per document

    async def extract_metadata_batch(
        self,
        file_previews: list[dict],
    ) -> list[dict]:
        """
        Extract metadata for multiple documents using LLM.

        Args:
            file_previews: List of dicts with 'filename' and 'content' (preview text)

        Returns:
            List of metadata dicts for each document
        """
        if not file_previews:
            return []

        # Build prompt for LLM
        prompt = self._build_extraction_prompt(file_previews)

        try:
            # Call LLM
            messages = [
                {
                    "role": "system",
                    "content": "你是文档元数据提取专家。请严格按照用户要求的JSON格式返回结果，不要返回其他任何内容。",
                },
                {"role": "user", "content": prompt},
            ]

            response = await llm_service.chat(
                messages=messages,
                temperature=0.3,  # Lower temperature for more consistent output
                max_tokens=4096,
            )

            # Parse LLM response
            results = self._parse_llm_response(response, len(file_previews))

            # Fill in defaults for failed extractions
            for i, result in enumerate(results):
                if not result or not result.get("doc_name"):
                    results[i] = self._get_default_metadata(
                        file_previews[i]["filename"] if i < len(file_previews) else "未命名文档"
                    )
                results[i]["file_index"] = i
                results[i]["filename"] = file_previews[i]["filename"] if i < len(file_previews) else ""

            return results

        except Exception as e:
            logger.error(f"LLM metadata extraction failed: {e}")
            # Return defaults for all documents
            return [
                {
                    **self._get_default_metadata(fp["filename"]),
                    "file_index": i,
                    "filename": fp["filename"],
                }
                for i, fp in enumerate(file_previews)
            ]

    def _build_extraction_prompt(self, file_previews: list[dict]) -> str:
        """Build the prompt for LLM metadata extraction"""
        docs_text = []
        for i, fp in enumerate(file_previews):
            preview = fp.get("content", "")[: self.max_preview_length]
            docs_text.append(f"[文档{i}] 文件名: {fp['filename']}\n内容预览:\n{preview}\n---")

        prompt = f"""请分析以下{len(file_previews)}个文档的内容，为每个文档提取元数据。

要求返回JSON数组，每个元素包含以下字段：
- doc_name: 文档标题（简洁准确，不超过50字，如果内容中有明确标题则使用，否则根据内容概括）
- department: 所属部门（如技术部、安全部、人事部、财务部、运营部等，如果无法判断则留空""）
- category: 文档类别（如管理制度、操作规程、培训教材、应急预案、技术文档、会议纪要等，如果无法判断则留空""）
- security_level: 密级（公开/内部/机密，根据内容敏感程度判断，默认"内部"）
- tags: 标签列表（3-5个关键词，用于描述文档主题和内容）
- description: 文档描述（100字以内，概括文档的主要内容和用途）

文档列表：
{chr(10).join(docs_text)}

请仅返回JSON数组，不要包含任何其他文字或代码块标记："""

        return prompt

    def _parse_llm_response(self, response: str, expected_count: int) -> list[dict]:
        """Parse LLM response to extract metadata list"""
        try:
            # Clean response - remove markdown code blocks if present
            cleaned = response.strip()
            if cleaned.startswith("```"):
                # Remove code block markers
                lines = cleaned.split("\n")
                if lines[0].startswith("```"):
                    lines = lines[1:]
                if lines and lines[-1].strip() == "```":
                    lines = lines[:-1]
                cleaned = "\n".join(lines)

            # Try to find JSON array in the response
            # Look for [ ... ] pattern
            match = re.search(r'\[[\s\S]*\]', cleaned)
            if match:
                cleaned = match.group(0)

            # Parse JSON
            results = json.loads(cleaned)

            if isinstance(results, list):
                return results
            elif isinstance(results, dict):
                # Single result wrapped
                return [results]
            else:
                return []

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse LLM response as JSON: {e}")
            logger.debug(f"Response was: {response[:500]}")
            return []
        except Exception as e:
            logger.error(f"Error parsing LLM response: {e}")
            return []

    def _get_default_metadata(self, filename: str) -> dict:
        """Get default metadata when extraction fails"""
        # Remove extension from filename for doc_name
        doc_name = filename
        if "." in filename:
            doc_name = filename.rsplit(".", 1)[0]

        return {
            "doc_name": doc_name,
            "department": "",
            "category": "",
            "security_level": "内部",
            "tags": [],
            "description": "",
        }

    async def extract_text_preview(
        self,
        content: bytes,
        file_type: str,
        max_length: int = 2000,
    ) -> str:
        """
        Extract text preview from file content for metadata extraction.
        Reuses parsing logic from document processor.
        """
        try:
            if file_type in ["txt", "md", "csv"]:
                # Plain text files
                try:
                    text = content.decode("utf-8")
                except UnicodeDecodeError:
                    try:
                        text = content.decode("gbk")
                    except UnicodeDecodeError:
                        text = content.decode("utf-8", errors="replace")
                return text[:max_length]

            elif file_type == "pdf":
                return await self._extract_pdf_preview(content, max_length)

            elif file_type == "docx":
                return await self._extract_docx_preview(content, max_length)

            elif file_type == "doc":
                # Check if it's actually a docx (ZIP format)
                if content[:4] == b'PK\x03\x04':
                    return await self._extract_docx_preview(content, max_length)
                return await self._extract_doc_preview(content, max_length)

            elif file_type in ["xls", "xlsx"]:
                return await self._extract_excel_preview(content, max_length)

            elif file_type in ["ppt", "pptx"]:
                return await self._extract_ppt_preview(content, max_length)

            else:
                return f"[{file_type.upper()}文件]"

        except Exception as e:
            logger.error(f"Error extracting preview for {file_type}: {e}")
            return ""

    async def _extract_pdf_preview(self, content: bytes, max_length: int) -> str:
        """Extract text preview from PDF"""
        try:
            import pdfplumber

            text_parts = []
            total_len = 0

            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages[:5]:  # Only first 5 pages for preview
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                        total_len += len(text)
                        if total_len >= max_length:
                            break

            return "\n\n".join(text_parts)[:max_length]

        except ImportError:
            return "[PDF文件]"
        except Exception as e:
            logger.error(f"Error extracting PDF preview: {e}")
            return ""

    async def _extract_docx_preview(self, content: bytes, max_length: int) -> str:
        """Extract text preview from DOCX"""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            text_parts = []
            total_len = 0

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
                    total_len += len(para.text)
                    if total_len >= max_length:
                        break

            return "\n\n".join(text_parts)[:max_length]

        except ImportError:
            return "[Word文件]"
        except Exception as e:
            logger.error(f"Error extracting DOCX preview: {e}")
            return ""

    async def _extract_doc_preview(self, content: bytes, max_length: int) -> str:
        """Extract text preview from legacy DOC format"""
        try:
            import olefile

            ole = olefile.OleFileIO(io.BytesIO(content))
            if ole.exists("WordDocument"):
                word_stream = ole.openstream("WordDocument").read()
                # Try to extract readable text
                text = word_stream.decode("utf-16-le", errors="ignore")
                # Clean up
                import re

                lines = []
                for segment in re.split(r'[\x00-\x08\x0b\x0c\x0e-\x1f]+', text):
                    cleaned = segment.strip()
                    if cleaned and len(cleaned) > 1:
                        lines.append(cleaned)
                ole.close()
                return "\n".join(lines)[:max_length]
            ole.close()
        except ImportError:
            pass
        except Exception as e:
            logger.debug(f"olefile extraction failed: {e}")

        return "[DOC文件]"

    async def _extract_excel_preview(self, content: bytes, max_length: int) -> str:
        """Extract text preview from Excel"""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text_parts = []
            total_len = 0

            for sheet_name in wb.sheetnames[:3]:  # First 3 sheets
                ws = wb[sheet_name]
                text_parts.append(f"=== {sheet_name} ===")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(cell) if cell is not None else "" for cell in row]
                    if any(v for v in row_vals):
                        row_text = "\t".join(row_vals)
                        text_parts.append(row_text)
                        total_len += len(row_text)
                        row_count += 1
                        if row_count >= 20 or total_len >= max_length:
                            break
                if total_len >= max_length:
                    break

            wb.close()
            return "\n".join(text_parts)[:max_length]

        except ImportError:
            return "[Excel文件]"
        except Exception as e:
            logger.error(f"Error extracting Excel preview: {e}")
            return ""

    async def _extract_ppt_preview(self, content: bytes, max_length: int) -> str:
        """Extract text preview from PowerPoint"""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []
            total_len = 0

            for i, slide in enumerate(prs.slides[:10]):  # First 10 slides
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text_parts.append(f"[幻灯片{i+1}] " + " ".join(slide_texts))
                    total_len += len(text_parts[-1])
                    if total_len >= max_length:
                        break

            return "\n\n".join(text_parts)[:max_length]

        except ImportError:
            return "[PPT文件]"
        except Exception as e:
            logger.error(f"Error extracting PPT preview: {e}")
            return ""


# Singleton instance
metadata_extractor = MetadataExtractor()
