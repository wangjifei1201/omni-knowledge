import io
import json

from core.database import get_db
from core.security import get_current_admin, get_current_user, get_current_user_optional
from fastapi import (
    APIRouter,
    BackgroundTasks,
    Depends,
    File,
    Form,
    HTTPException,
    Query,
    UploadFile,
)
from fastapi.responses import StreamingResponse
from models.document import Chapter, Document, DocumentChunk, Tag
from models.user import User
from schemas.document import (
    ChapterResponse,
    ChunkingConfigUpdate,
    ChunkListResponse,
    ChunkResponse,
    DocumentCreate,
    DocumentFilter,
    DocumentListResponse,
    DocumentResponse,
    DocumentUpdate,
)
from services.document.batch_train_manager import batch_train_manager
from services.document.chunking_strategy import chunking_engine
from services.document.metadata_extractor import metadata_extractor
from services.document.processor import document_processor
from services.rag.vector_store import faiss_vector_store
from services.storage.local_storage import local_storage
from sqlalchemy import asc, delete, desc, func, or_, select
from sqlalchemy.ext.asyncio import AsyncSession

router = APIRouter(prefix="/documents", tags=["文档管理"])

ALLOWED_EXTENSIONS = {
    "pdf",
    "doc",
    "docx",
    "xls",
    "xlsx",
    "ppt",
    "pptx",
    "txt",
    "md",
    "csv",
    "png",
    "jpg",
    "jpeg",
    "gif",
    "bmp",
    "mp3",
    "wav",
    "mp4",
}


def get_file_extension(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


@router.post("", response_model=DocumentResponse)
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    doc_name: str = Form(""),
    department: str = Form(""),
    category: str = Form(""),
    security_level: str = Form("内部"),
    tags: str = Form(""),  # comma-separated
    description: str = Form(""),
    effective_date: str = Form(""),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    ext = get_file_extension(file.filename or "")
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext}")

    file_content = await file.read()
    file_size = len(file_content)
    final_name = doc_name or file.filename or "未命名文档"

    tag_list = [t.strip() for t in tags.split(",") if t.strip()] if tags else []

    # Create document record first to get ID
    doc = Document(
        doc_name=final_name,
        file_path="",  # Will update after storage
        file_type=ext,
        file_size=file_size,
        department=department,
        category=category,
        security_level=security_level,
        description=description,
        effective_date=effective_date,
        tags=json.dumps(tag_list, ensure_ascii=False),  # Store as JSON string
        uploaded_by=current_user.id,
        status="processing",
    )
    db.add(doc)
    await db.flush()
    await db.refresh(doc)

    # Upload to local storage
    storage_info = await local_storage.upload_file(
        file_content=file_content,
        original_name=file.filename or final_name,
        doc_id=doc.id,
        category="documents",
    )
    doc.file_path = storage_info["storage_path"]
    await db.flush()

    # Add tags
    for tag_name in tag_list:
        db.add(Tag(doc_id=doc.id, tag_name=tag_name))

    # Schedule document processing in background with default strategy
    background_tasks.add_task(
        process_document_task,
        doc.id,
        doc.file_path,
        doc.file_type,
        doc.doc_name,
        doc.chunking_strategy,
        doc.chunking_params,
    )

    return DocumentResponse.model_validate(doc)


@router.post("/extract-metadata")
async def extract_document_metadata(
    files: list[UploadFile] = File(...),
    current_user: User = Depends(get_current_user),
):
    """
    Extract metadata from multiple documents using LLM.
    Max 10 files per request.
    """
    if len(files) > 10:
        raise HTTPException(status_code=400, detail="一次最多上传10个文件")

    if len(files) == 0:
        raise HTTPException(status_code=400, detail="请选择要上传的文件")

    # Validate file types and read content
    file_previews = []
    for file in files:
        ext = get_file_extension(file.filename or "")
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext} ({file.filename})")

        # Read file content
        content = await file.read()
        await file.seek(0)  # Reset file position

        # Extract text preview
        preview = await metadata_extractor.extract_text_preview(content, ext)

        file_previews.append(
            {
                "filename": file.filename or "未命名文件",
                "content": preview,
            }
        )

    # Extract metadata using LLM
    results = await metadata_extractor.extract_metadata_batch(file_previews)

    return {"results": results}


@router.post("/batch-train")
async def batch_train_documents(
    background_tasks: BackgroundTasks,
    request_body: dict,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """
    Start batch training for multiple documents.
    Re-processes documents and regenerates vector embeddings.
    """
    doc_ids = request_body.get("doc_ids", [])

    if not doc_ids:
        raise HTTPException(status_code=400, detail="请提供要训练的文档ID列表")

    if len(doc_ids) == 0:
        raise HTTPException(status_code=400, detail="请选择要训练的文档")

    if len(doc_ids) > 50:
        raise HTTPException(status_code=400, detail="一次最多训练50个文档")

    # Verify all documents exist and get their names
    doc_names = {}
    for doc_id in doc_ids:
        result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
        doc = result.scalar_one_or_none()
        if not doc:
            raise HTTPException(status_code=404, detail=f"文档不存在: {doc_id}")
        if doc.status == "processing":
            raise HTTPException(status_code=400, detail=f"文档正在处理中: {doc.doc_name}")
        doc_names[doc_id] = doc.doc_name

    # Create batch train task
    task_id = batch_train_manager.create_task(doc_ids, doc_names)

    # Schedule background execution
    background_tasks.add_task(batch_train_manager.execute_batch_train, task_id)

    return {
        "task_id": task_id,
        "total": len(doc_ids),
        "message": "批量训练任务已创建",
    }


@router.get("/batch-train/{task_id}/status")
async def get_batch_train_status(
    task_id: str,
    current_user: User = Depends(get_current_user),
):
    """
    Get the status of a batch training task.
    """
    status = batch_train_manager.get_task_status(task_id)
    if not status:
        raise HTTPException(status_code=404, detail="任务不存在")
    return status


@router.get("/chunking-strategies")
async def get_chunking_strategies(
    current_user: User = Depends(get_current_user),
):
    """Return available chunking strategy definitions"""
    return {"strategies": chunking_engine.get_strategy_definitions()}


async def process_document_task(
    doc_id: str,
    file_path: str,
    file_type: str,
    doc_name: str,
    chunking_strategy: str = "paragraph",
    chunking_params: str | None = None,
):
    """Background task to process document"""
    from core.database import AsyncSessionLocal

    # Parse chunking_params from JSON string
    params_dict = None
    if chunking_params:
        try:
            params_dict = json.loads(chunking_params) if isinstance(chunking_params, str) else chunking_params
        except (json.JSONDecodeError, TypeError):
            params_dict = None

    result = await document_processor.process_document(
        doc_id=doc_id,
        file_path=file_path,
        file_type=file_type,
        doc_name=doc_name,
        chunking_strategy=chunking_strategy,
        chunking_params=params_dict,
    )

    # Update document status in database
    async with AsyncSessionLocal() as db:
        doc_result = await db.execute(select(Document).where(Document.id == doc_id))
        doc = doc_result.scalar_one_or_none()
        if doc:
            doc.status = result.get("status", "failed")
            doc.chunk_count = result.get("chunk_count", 0)
            doc.page_count = result.get("page_count", 0)
            if result.get("error"):
                doc.error_message = result["error"]
            await db.commit()


@router.get("", response_model=DocumentListResponse)
async def list_documents(
    keyword: str = Query("", description="搜索关键词"),
    department: str = Query("", description="部门筛选"),
    category: str = Query("", description="类别筛选"),
    file_type: str = Query("", description="文件类型"),
    status: str = Query("", description="处理状态"),
    page: int = Query(1, ge=1),
    page_size: int = Query(20, ge=1, le=100),
    sort_by: str = Query("created_at"),
    sort_order: str = Query("desc"),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    query = select(Document).where(Document.is_deleted == False)

    if keyword:
        query = query.where(
            or_(
                Document.doc_name.ilike(f"%{keyword}%"),
                Document.description.ilike(f"%{keyword}%"),
            )
        )
    if department:
        query = query.where(Document.department == department)
    if category:
        query = query.where(Document.category == category)
    if file_type:
        query = query.where(Document.file_type == file_type)
    if status:
        query = query.where(Document.status == status)

    # Count total
    count_query = select(func.count()).select_from(query.subquery())
    total_result = await db.execute(count_query)
    total = total_result.scalar() or 0

    # Sort
    sort_column = getattr(Document, sort_by, Document.created_at)
    if sort_order == "asc":
        query = query.order_by(asc(sort_column))
    else:
        query = query.order_by(desc(sort_column))

    # Paginate
    query = query.offset((page - 1) * page_size).limit(page_size)
    result = await db.execute(query)
    docs = result.scalars().all()

    return DocumentListResponse(
        items=[DocumentResponse.model_validate(d) for d in docs],
        total=total,
        page=page,
        page_size=page_size,
    )


@router.get("/stats/overview")
async def get_document_stats(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    total = await db.execute(select(func.count()).where(Document.is_deleted == False))
    by_dept = await db.execute(
        select(Document.department, func.count()).where(Document.is_deleted == False).group_by(Document.department)
    )
    by_type = await db.execute(
        select(Document.file_type, func.count()).where(Document.is_deleted == False).group_by(Document.file_type)
    )
    by_status = await db.execute(
        select(Document.status, func.count()).where(Document.is_deleted == False).group_by(Document.status)
    )
    total_size = await db.execute(select(func.sum(Document.file_size)).where(Document.is_deleted == False))

    return {
        "total_documents": total.scalar() or 0,
        "total_size": total_size.scalar() or 0,
        "by_department": dict(by_dept.all()),
        "by_file_type": dict(by_type.all()),
        "by_status": dict(by_status.all()),
    }


@router.get("/{doc_id}", response_model=DocumentResponse)
async def get_document(
    doc_id: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")
    return DocumentResponse.model_validate(doc)


@router.put("/{doc_id}", response_model=DocumentResponse)
async def update_document(
    doc_id: str,
    req: DocumentUpdate,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    update_data = req.model_dump(exclude_unset=True)
    # Convert tags list to JSON string
    if "tags" in update_data and update_data["tags"] is not None:
        update_data["tags"] = json.dumps(update_data["tags"], ensure_ascii=False)

    for key, value in update_data.items():
        setattr(doc, key, value)
    await db.flush()
    await db.refresh(doc)
    return DocumentResponse.model_validate(doc)


@router.put("/{doc_id}/chunking-config", response_model=DocumentResponse)
async def update_chunking_config(
    doc_id: str,
    req: ChunkingConfigUpdate,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Update document chunking strategy configuration (save only, no retrain)"""
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    # Validate strategy name
    valid_strategies = {"character", "paragraph", "heading"}
    if req.chunking_strategy not in valid_strategies:
        raise HTTPException(status_code=400, detail=f"无效的分段策略: {req.chunking_strategy}")

    doc.chunking_strategy = req.chunking_strategy
    doc.chunking_params = json.dumps(req.chunking_params, ensure_ascii=False) if req.chunking_params else None
    await db.flush()
    await db.refresh(doc)
    return DocumentResponse.model_validate(doc)


@router.post("/{doc_id}/retrain")
async def retrain_document(
    doc_id: str,
    background_tasks: BackgroundTasks,
    request_body: dict | None = None,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Save chunking config (if provided) and retrain document"""
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    if not doc.file_path:
        raise HTTPException(status_code=400, detail="文件不存在，无法重新训练")

    if doc.status == "processing":
        raise HTTPException(status_code=400, detail="文档正在处理中，请等待完成")

    # Update chunking config if provided
    if request_body:
        strategy = request_body.get("chunking_strategy")
        params = request_body.get("chunking_params")
        if strategy:
            valid_strategies = {"character", "paragraph", "heading"}
            if strategy not in valid_strategies:
                raise HTTPException(status_code=400, detail=f"无效的分段策略: {strategy}")
            doc.chunking_strategy = strategy
        if params is not None:
            doc.chunking_params = json.dumps(params, ensure_ascii=False) if params else None

    # Delete existing vectors
    await faiss_vector_store.delete_by_doc_id(doc_id)

    # Update status to processing
    doc.status = "processing"
    doc.parse_progress = 0.0
    await db.flush()

    # Capture values before session closes
    file_path = doc.file_path
    file_type = doc.file_type
    doc_name = doc.doc_name
    chunking_strategy = doc.chunking_strategy
    chunking_params = doc.chunking_params

    # Schedule reprocessing
    background_tasks.add_task(
        process_document_task,
        doc_id,
        file_path,
        file_type,
        doc_name,
        chunking_strategy,
        chunking_params,
    )

    return {"message": "已开始重新训练", "doc_id": doc_id}


@router.delete("/{doc_id}")
async def delete_document(
    doc_id: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    doc.is_deleted = True
    await db.flush()

    # Delete from local storage
    await local_storage.delete_doc_files(doc_id)

    # Delete from FAISS vector store
    await faiss_vector_store.delete_by_doc_id(doc_id)

    # Delete chunks from database
    await db.execute(delete(DocumentChunk).where(DocumentChunk.doc_id == doc_id))

    return {"message": "文档已删除"}


@router.get("/{doc_id}/download")
async def download_document(
    doc_id: str,
    current_user: User = Depends(get_current_user_optional),
    db: AsyncSession = Depends(get_db),
):
    """Download original document file"""
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    if not doc.file_path:
        raise HTTPException(status_code=404, detail="文件不存在")

    file_exists = await local_storage.file_exists(doc.file_path)
    if not file_exists:
        raise HTTPException(status_code=404, detail="文件不存在")

    # Capture file path before db session closes
    file_path = doc.file_path
    doc_name = doc.doc_name
    file_type = doc.file_type

    # Get file info for content-length header
    file_info = await local_storage.get_file_info(file_path)

    # Determine content type
    content_type_map = {
        "pdf": "application/pdf",
        "doc": "application/msword",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xls": "application/vnd.ms-excel",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "ppt": "application/vnd.ms-powerpoint",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "txt": "text/plain",
        "md": "text/markdown",
        "csv": "text/csv",
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
    }
    content_type = content_type_map.get(file_type, "application/octet-stream")

    # Read file content directly to avoid async generator issues with db session
    file_content = await local_storage.download_file(file_path)
    if file_content is None:
        raise HTTPException(status_code=404, detail="文件读取失败")

    # URL encode filename for Content-Disposition header
    from urllib.parse import quote

    from fastapi.responses import Response

    encoded_filename = quote(f"{doc_name}.{file_type}")

    return Response(
        content=file_content,
        media_type=content_type,
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{encoded_filename}",
            "Content-Length": str(len(file_content)),
        },
    )


@router.get("/{doc_id}/preview")
async def preview_document(
    doc_id: str,
    current_user: User = Depends(get_current_user_optional),
    db: AsyncSession = Depends(get_db),
):
    """预览文档 - 二进制格式（doc/docx/pdf/xlsx等）自动提取为纯文本返回"""
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    if not doc.file_path:
        raise HTTPException(status_code=404, detail="文件不存在")

    file_exists = await local_storage.file_exists(doc.file_path)
    if not file_exists:
        raise HTTPException(status_code=404, detail="文件不存在")

    # 在 db session 关闭前获取值
    file_path = doc.file_path
    doc_name = doc.doc_name
    file_type = doc.file_type

    # 读取文件内容
    file_content = await local_storage.download_file(file_path)
    if file_content is None:
        raise HTTPException(status_code=404, detail="文件读取失败")

    from urllib.parse import quote

    from fastapi.responses import Response

    # 对于二进制文档格式，提取纯文本后返回
    text_extractable_types = {"doc", "docx", "pdf", "xls", "xlsx", "ppt", "pptx"}
    # 纯文本格式直接返回
    plain_text_types = {"txt", "md", "csv", "json", "xml", "yaml", "yml", "html", "css", "js", "ts", "py"}

    if file_type in text_extractable_types:
        try:
            extracted_text = _extract_text_from_binary(file_content, file_type)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"文档内容提取失败: {str(e)}")

        return Response(
            content=extracted_text.encode("utf-8"),
            media_type="text/plain; charset=utf-8",
            headers={
                "Content-Disposition": f"inline; filename*=UTF-8''{quote(f'{doc_name}.txt')}",
            },
        )
    elif file_type in plain_text_types:
        # 纯文本尝试 utf-8 解码，失败则尝试 gbk
        try:
            text = file_content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = file_content.decode("gbk")
            except UnicodeDecodeError:
                text = file_content.decode("utf-8", errors="replace")

        return Response(
            content=text.encode("utf-8"),
            media_type="text/plain; charset=utf-8",
            headers={
                "Content-Disposition": f"inline; filename*=UTF-8''{quote(f'{doc_name}.{file_type}')}",
            },
        )
    else:
        # 图片等其他格式，直接返回原始内容
        content_type_map = {
            "png": "image/png",
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
            "gif": "image/gif",
            "bmp": "image/bmp",
        }
        content_type = content_type_map.get(file_type, "application/octet-stream")
        encoded_filename = quote(f"{doc_name}.{file_type}")

        return Response(
            content=file_content,
            media_type=content_type,
            headers={
                "Content-Disposition": f"inline; filename*=UTF-8''{encoded_filename}",
                "Content-Length": str(len(file_content)),
            },
        )


def _extract_text_from_doc(content: bytes) -> str:
    """从旧版 .doc (OLE Compound Document) 格式中提取纯文本"""
    import os
    import subprocess
    import tempfile

    # 方案1：尝试使用 textutil (macOS 自带)
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", tmp_path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        finally:
            os.unlink(tmp_path)
    except FileNotFoundError:
        pass  # textutil 不可用（非 macOS），继续尝试其他方案
    except Exception:
        pass

    # 方案2：尝试使用 antiword
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            result = subprocess.run(
                ["antiword", tmp_path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        finally:
            os.unlink(tmp_path)
    except FileNotFoundError:
        pass
    except Exception:
        pass

    # 方案3：使用 olefile 从 OLE 文档中提取原始文本
    try:
        import olefile

        buf = io.BytesIO(content)
        ole = olefile.OleFileIO(buf)
        # Word 文档的文本流
        if ole.exists("WordDocument"):
            # 尝试从 Word Document 流中提取
            word_stream = ole.openstream("WordDocument").read()
            # 提取可读文本（简化方式：过滤可打印 Unicode 字符）
            text = word_stream.decode("utf-16-le", errors="ignore")
            # 清理不可见字符，只保留可打印内容
            import re

            lines = []
            for segment in re.split(r'[\x00-\x08\x0b\x0c\x0e-\x1f]+', text):
                cleaned = segment.strip()
                if cleaned and len(cleaned) > 1:
                    lines.append(cleaned)
            if lines:
                ole.close()
                return "\n".join(lines)
        ole.close()
    except ImportError:
        pass
    except Exception:
        pass

    return "（无法提取 .doc 格式内容，请将文档转换为 .docx 格式后重新上传）"


def _extract_text_from_docx(content: bytes) -> str:
    """从 .docx 格式中提取纯文本"""
    from docx import Document as DocxDocument

    buf = io.BytesIO(content)
    doc = DocxDocument(buf)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    # 同时提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                paragraphs.append(row_text)
    return "\n\n".join(paragraphs) if paragraphs else "（文档内容为空）"


def _extract_text_from_binary(content: bytes, file_type: str) -> str:
    """从二进制文档中提取纯文本内容"""

    if file_type == "doc":
        # 检测实际文件格式：有些 .doc 文件实际是 .docx（ZIP格式）
        if content[:4] == b'PK\x03\x04':
            return _extract_text_from_docx(content)
        return _extract_text_from_doc(content)

    elif file_type == "docx":
        return _extract_text_from_docx(content)

    elif file_type == "pdf":
        import pdfplumber

        buf = io.BytesIO(content)
        pages_text = []
        with pdfplumber.open(buf) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    pages_text.append(f"--- 第 {i + 1} 页 ---\n{text}")
        return "\n\n".join(pages_text) if pages_text else "（PDF 内容为空）"

    elif file_type in ("xls", "xlsx"):
        from openpyxl import load_workbook

        buf = io.BytesIO(content)
        wb = load_workbook(buf, read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(cell) if cell is not None else "" for cell in row]
                if any(v for v in row_vals):
                    rows.append("\t".join(row_vals))
            if rows:
                sheets_text.append(f"=== 工作表: {sheet_name} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets_text) if sheets_text else "（表格内容为空）"

    elif file_type in ("ppt", "pptx"):
        try:
            from pptx import Presentation

            buf = io.BytesIO(content)
            prs = Presentation(buf)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides_text.append(f"--- 幻灯片 {i + 1} ---\n" + "\n".join(texts))
            return "\n\n".join(slides_text) if slides_text else "（演示文稿内容为空）"
        except ImportError:
            return "（不支持 PPT 格式预览，请安装 python-pptx）"

    return "（不支持的文档格式）"


@router.post("/{doc_id}/reparse")
async def reparse_document(
    doc_id: str,
    background_tasks: BackgroundTasks,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Re-trigger document parsing"""
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    if not doc.file_path:
        raise HTTPException(status_code=400, detail="文件不存在，无法重新解析")

    # Delete existing vectors for this document
    await faiss_vector_store.delete_by_doc_id(doc_id)

    # Update status
    doc.status = "processing"
    await db.flush()

    # Capture values
    file_path = doc.file_path
    file_type = doc.file_type
    doc_name = doc.doc_name
    chunking_strategy = doc.chunking_strategy
    chunking_params = doc.chunking_params

    # Schedule processing
    background_tasks.add_task(
        process_document_task,
        doc.id,
        file_path,
        file_type,
        doc_name,
        chunking_strategy,
        chunking_params,
    )

    return {"message": "已开始重新解析", "doc_id": doc_id}


@router.get("/{doc_id}/chunks", response_model=ChunkListResponse)
async def get_document_chunks(
    doc_id: str,
    page: int = Query(1, ge=1),
    page_size: int = Query(20, ge=1, le=100),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Get document chunks from database with pagination"""
    result = await db.execute(select(Document).where(Document.id == doc_id, Document.is_deleted == False))
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="文档不存在")

    # Count total chunks
    count_result = await db.execute(select(func.count()).where(DocumentChunk.doc_id == doc_id))
    total = count_result.scalar() or 0

    # Query chunks with pagination
    chunks_result = await db.execute(
        select(DocumentChunk)
        .where(DocumentChunk.doc_id == doc_id)
        .order_by(DocumentChunk.chunk_index)
        .offset((page - 1) * page_size)
        .limit(page_size)
    )
    chunks = chunks_result.scalars().all()

    return ChunkListResponse(
        items=[ChunkResponse.model_validate(c) for c in chunks],
        total=total,
        page=page,
        page_size=page_size,
    )


@router.get("/{doc_id}/chapters", response_model=list[ChapterResponse])
async def get_chapters(
    doc_id: str,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Chapter).where(Chapter.doc_id == doc_id).order_by(Chapter.order_index))
    chapters = result.scalars().all()
    return [ChapterResponse.model_validate(c) for c in chapters]
